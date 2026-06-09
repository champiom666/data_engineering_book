"""Build a 100-page .pptx from the outline markdown, with real book images.

Input:  outputs/presentations/2026-05-28-1hour-100page-outline.md
Output: outputs/presentations/2026-05-28-1hour-100page-presentation.pptx

Images sourced from docs/images/partN/ — each page is mapped to the most relevant
book figure by chapter + keyword. Pages without a good match get a styled
text-only diagram block.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


REPO = Path(__file__).resolve().parent.parent
INPUT_MD = REPO / "outputs/presentations/2026-05-28-1hour-100page-outline.md"
OUTPUT_PPTX = REPO / "outputs/presentations/2026-05-28-1hour-100page-presentation.pptx"
IMG_ROOT = REPO / "docs/images"
STRUCTURE_IMG = REPO / "images/structure_cn.png"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Color palette — Academic White (Option A)
C_BG = RGBColor(0xFF, 0xFF, 0xFF)        # pure white
C_PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)   # deep indigo (titles / module tag)
C_ACCENT = RGBColor(0xEA, 0x58, 0x0C)    # warm orange (accent strip / underline)
C_INK = RGBColor(0x1F, 0x29, 0x37)        # near-black slate
C_MUTED = RGBColor(0x6B, 0x72, 0x80)      # muted grey
C_BOX_BG = RGBColor(0xF9, 0xFA, 0xFB)     # soft grey card background
C_BOX_BORDER = RGBColor(0xCB, 0xD5, 0xE1) # cool grey border
C_TAKEAWAY_BG = RGBColor(0x1E, 0x3A, 0x8A)  # deep indigo strip
C_TAKEAWAY_TEXT = RGBColor(0xFF, 0xFF, 0xFF)


# Page-type → bottom-label mapping (Option B for the takeaway band).
# Pages that don't match fall back to "核心论断".
TAKEAWAY_LABEL_BY_TYPE: dict[str, str] = {
    "痛点": "代价",
    "警示": "代价",
    "关键技术": "适用场景",
    "方法地图": "全流程",
    "一图速览": "金句",
    "海报": "金句",
    "收尾": "金句",
    "实战场景": "关键数据",
    "架构图": "关键数据",
    "关键挑战": "关键数据",
    "成果": "关键数据",
    "模板": "上手第一步",
    "Checklist": "上手第一步",
    "概念": "核心论断",
    "组织": "核心论断",
    "技术地图": "核心论断",
    "总结": "核心论断",
    "过渡": "核心论断",
    "趋势": "核心论断",
    "工具": "核心论断",
    "路径": "核心论断",
    "矩阵": "核心论断",
    "资源": "核心论断",
    "封面": "金句",
    "自我介绍": "今天的承诺",
}


def takeaway_label_for(page_type: str) -> str:
    """Match the page_type string against the mapping. Loose substring match."""
    if not page_type:
        return "核心论断"
    # Page type field may carry compound descriptors like "实战场景 / 总览"
    # Try each segment in order.
    for seg in re.split(r"[/／、,，\s]+", page_type):
        seg = seg.strip()
        if not seg:
            continue
        if seg in TAKEAWAY_LABEL_BY_TYPE:
            return TAKEAWAY_LABEL_BY_TYPE[seg]
    return "核心论断"


# =====================================================================
# Page → image mapping. Keys are page numbers, values are paths relative
# to docs/images/. Pages NOT in this dict will render a text-only diagram.
# Curated by page topic + book chapter.
# =====================================================================

PAGE_IMAGES: dict[int, str] = {
    # Module 0 — opening (no real images for cover/poster slides)
    2: "../../images/structure_cn.png",  # commitments backed by full book map
    3: "part1/data_lifecycle_map_1775830407042.png",  # failure project intro
    5: "../../images/structure_cn.png",  # full-book structure
    # Module 1 — paradigm & lifecycle (Ch01-03)
    7: "part1/data_engineering_roles_1775830393574.png",
    8: "part1/data_quality_hierarchy_1775835516841.png",
    9: "part1/defect_metric_radar_1775835533937.png",
    10: "part1/data_engineering_roles_1775830393574.png",
    11: "part2/data_operations_flywheel.png",
    12: "part1/data_lifecycle_map_1775830407042.png",
    13: "part1/data_quality_hierarchy_1775835516841.png",
    14: "part1/data_quality_hierarchy_1775835516841.png",
    15: "part1/data_quality_gates_1775835548587.png",
    16: "part1/ai_data_stack_architecture.png",
    17: "part1/ai_data_stack_architecture.png",
    18: "part1/cost_governance_loop.png",
    19: "../../images/structure_cn.png",
    20: "../../images/structure_cn.png",
    # L1 Text pretraining (p21-28, Ch04-07)
    21: "part2/pretrain_data_source_map.png",
    22: "part2/cleaning_pipeline_overview.png",
    23: "part2/cleaning_pipeline_overview.png",
    24: "part2/quality_filter_funnel_loop.png",
    25: "part2/quality_filter_funnel_loop.png",
    26: "part2/data_ingestion_provenance_chain.png",
    27: "part2/data_evaluation_loop.png",
    28: "part2/quality_filter_funnel_loop.png",
    # L2 Multimodal (p29-35, Ch08-11)
    29: "part3/multimodal_data_panorama.png",
    30: "part3/multimodal_data_panorama.png",
    31: "part3/image_semantic_alignment_flow.png",
    32: "part3/recaptioning_ocr_pipeline.png",
    33: "part3/av_sample_pipeline.png",
    34: "part3/cross_modal_alignment_hierarchy.png",
    35: "part3/fusion_training_sample_design.png",
    # L3 SFT/Preference (p36-41, Ch12-14)
    36: "part4/sft_instruction_architecture.png",
    37: "part4/sft_generation_loop.png",
    38: "part4/sft_instruction_architecture.png",
    39: "part4/图11_1_人类偏好示意图.png",
    40: "part4/图14_1.png",
    41: "part4/图13_1.png",
    # L4 Synthetic (p42-47, Ch15-17)
    42: "part4/图9_1_自我指令和进化指令对比.png",
    43: "part5/图15_1.png",
    44: "part4/图9_1_自我指令和进化指令对比.png",
    45: "part5/图16_1.png",
    46: "part5/图17_1.png",
    47: "part5/图15_2.png",
    # L5 Reasoning / Agent (p48-53, Ch18-20)
    48: "part6/图18_1.png",
    49: "part6/图18_1.png",
    50: "part6/图18_2.png",
    51: "part6/图19_1.png",
    52: "part6/图20_1.png",
    53: "part6/图20_2.png",
    # L6 RAG / feedback (p54-59, Ch21-23)
    54: "part7/图21_1zh.png",
    55: "part7/图21_2zh.png",
    56: "part7/图21_3zh.png",
    57: "part7/图22_1zh.png",
    58: "part7/图23_1zh.png",
    59: "part7/图21_9zh.png",
    # L7 DataOps / platform (p60-65, Ch24-26)
    60: "part8/图24_1_LLM数据团队组织演进路径.png",
    61: "part8/图24_3_DataOps团队组织全景图.png",
    62: "part8/图24_2_DataOps飞轮四池协同示意图.png",
    63: "part8/图25_1_数据谱系与实验追踪图.png",
    64: "part8/图26_2_平台可观测性全景图.png",
    65: "part8/图25_2_版本管理体系全景图.png",
    # L8 Privacy / Compliance (p66-70, Ch27-28)
    66: "part9/图27_1_合规左移与治理协同架构图.png",
    67: "part9/图28_11_合规治理隐私流水线联邦训练与应用能力闭环图.png",
    68: "part9/图27_5_从数据接入到模型训练的合规门禁流程图.png",
    69: "part9/图28_10_联邦系统整体架构图.png",
    70: "part9/图28_4_隐私增强技术全景矩阵图.png",
    # Module 3 — three projects (P1 Mini-C4 p71-77, P2 LLaVA p78-84, P3 DataOps p85-90)
    71: "part10/10_1_fig01_mini_c4_pipeline_overview.png",
    72: "part10/10_1_fig01_mini_c4_pipeline_overview.png",
    73: "part2/cleaning_pipeline_overview.png",
    74: "part10/10_1_fig04_dedup_minhash_lsh.png",
    75: "part10/10_1_fig06_quality_filter.png",
    76: "part10/10_1_fig07_three_iterations.png",
    77: "part10/10_1_fig11_methodology_summary.png",
    78: "part10/10_3_fig01_llava_factory_overview.png",
    79: "part10/10_3_fig01_llava_factory_overview.png",
    80: "part10/10_3_fig05_bbox_alignment.png",
    81: "part10/10_5_fig06_multi_image_prompting.png",
    82: "part10/10_3_fig07_failure_attribution.png",
    83: "part10/10_5_fig08_optimization_roadmap.png",
    84: "part10/10_5_fig09_hybrid_rag.png",
    85: "part10/10_8_fig01_dataops_platform_overview.png",
    86: "part10/10_8_fig01_dataops_platform_overview.png",
    87: "part10/10_8_fig06_experiment_tracking.png",
    88: "part10/10_8_fig11_validation_pipeline.png",
    89: "part10/10_8_fig07_lineage_graph.png",
    90: "part10/10_10_fig05_architecture_code_mapping.png",
    # Module 4 — closing
    91: "../../images/structure_cn.png",
    92: "part10/10_10_fig01_flywheel_overview.png",
    93: "part5/图17_1.png",
    94: "part10/10_10_fig04_stage_plan.png",
    95: "../../images/structure_cn.png",
    96: "part8/图24_3_DataOps团队组织全景图.png",
    97: "part10/10_10_fig04_stage_plan.png",
    # 98, 99, 100 stay text-only (poster/resource/Q&A)
}


@dataclass
class Page:
    n: int
    title: str
    module: str = ""
    book_pos: str = ""
    page_type: str = ""
    points: list[str] = field(default_factory=list)
    speech: str = ""
    figure: str = ""
    takeaway: str = ""


# =====================================================================
# Outline parsing
# =====================================================================


def parse_outline(md_text: str) -> list[Page]:
    chunks = re.split(r"(?m)^### Page (\d+)\s*[—-]\s*(.+?)$", md_text)
    pages: list[Page] = []
    for i in range(1, len(chunks), 3):
        n = int(chunks[i])
        title = chunks[i + 1].strip()
        body = chunks[i + 2]
        pages.append(_parse_body(n, title, body))
    pages.sort(key=lambda p: p.n)
    return pages


def _extract_field(body: str, label: str) -> str:
    pattern = rf"\*\*{label}\*\*\s*[:：]\s*([^\n]+)"
    m = re.search(pattern, body)
    return m.group(1).strip() if m else ""


def _parse_body(n: int, title: str, body: str) -> Page:
    p = Page(n=n, title=title)
    p.module = _extract_field(body, "所属模块")
    p.book_pos = _extract_field(body, "对应书中位置")
    p.page_type = _extract_field(body, "页面类型")

    points_match = re.search(
        r"\*\*核心要点\*\*[^\n]*\n((?:\s{2,}-\s+[^\n]+\n?)+)", body
    )
    if points_match:
        for line in points_match.group(1).splitlines():
            m = re.match(r"\s+-\s+(.*\S)", line)
            if m:
                p.points.append(m.group(1).strip())

    speech_match = re.search(
        r"\*\*讲稿提示\*\*[^\n]*\n((?:\s*>\s?[^\n]*\n?)+)", body
    )
    if speech_match:
        lines = [
            re.sub(r"^\s*>\s?", "", ln).rstrip()
            for ln in speech_match.group(1).splitlines()
        ]
        p.speech = " ".join(line for line in lines if line).strip()

    p.figure = _extract_field(body, "配图建议")
    p.takeaway = _extract_field(body, "听众带走")
    p.takeaway = re.sub(r"[（(]一句话[）)]\s*$", "", p.takeaway).strip()
    return p


# =====================================================================
# Drawing helpers
# =====================================================================


def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = line_w
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=Pt(14),
    color=C_INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name="PingFang SC",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=Pt(16), color=C_INK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    for i, text in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(6)
        para.line_spacing = 1.25
        run = para.add_run()
        run.text = f"• {text}"
        run.font.size = size
        run.font.color.rgb = color
        run.font.name = "PingFang SC"
    return tb


def add_image_fitted(slide, image_path: Path, x, y, w, h):
    """Insert image into a (w, h) bounding box, centered, preserving aspect ratio."""
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        # Image is wider than box -> fit to width
        new_w = w
        new_h = int(w * ih / iw)
    else:
        new_h = h
        new_w = int(h * iw / ih)
    new_x = int(x + (w - new_w) / 2)
    new_y = int(y + (h - new_h) / 2)
    slide.shapes.add_picture(str(image_path), new_x, new_y, new_w, new_h)


def resolve_image_path(rel: str) -> Path | None:
    """Resolve image relative-path under docs/images/ or special prefix."""
    # Allow escape: "../../images/structure_cn.png" -> REPO/images/...
    if rel.startswith("../../images/"):
        p = REPO / rel.replace("../../", "")
    else:
        p = IMG_ROOT / rel
    return p if p.exists() else None


# =====================================================================
# Slide builders
# =====================================================================


def add_cover_slide(prs: Presentation, total_pages: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)
    add_rect(slide, 0, 0, Inches(0.35), SLIDE_H, C_ACCENT)

    add_text(
        slide,
        Inches(0.9),
        Inches(1.0),
        Inches(11.5),
        Inches(0.5),
        "DATA ENGINEERING FOR LLMs · 60-MIN EXECUTIVE OVERVIEW",
        size=Pt(14),
        color=C_PRIMARY,
        bold=True,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(1.7),
        Inches(11.5),
        Inches(1.8),
        "《大模型数据工程》",
        size=Pt(54),
        color=C_INK,
        bold=True,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(3.0),
        Inches(11.5),
        Inches(1.2),
        "架构、算法与项目实战 · 100 页全景导览",
        size=Pt(28),
        color=C_PRIMARY,
        bold=True,
    )
    add_text(
        slide,
        Inches(0.9),
        Inches(4.3),
        Inches(11.5),
        Inches(0.8),
        "1 小时 · 5 大模块 · 100 页 · 10 篇 28 章 + 10 实战",
        size=Pt(18),
        color=C_MUTED,
    )
    add_rect(slide, 0, Inches(7.0), SLIDE_W, Inches(0.5), C_TAKEAWAY_BG)
    add_text(
        slide,
        Inches(0.9),
        Inches(7.0),
        Inches(11.5),
        Inches(0.5),
        "2026-05-28   |   汇报人：__________   |   datascale-ai.github.io/data_engineering_book",
        size=Pt(12),
        color=C_TAKEAWAY_TEXT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _split_module_label(module: str) -> tuple[str, str]:
    """Split '模块 1 · 范式变革' -> ('模块 1', '范式变革'). Robust to no separator."""
    parts = re.split(r"\s*[·•]\s*", module, maxsplit=1)
    if len(parts) == 2:
        return parts[0].strip(), parts[1].strip()
    return module.strip(), ""


# Special-cased poster slides: render full-width without right-side image.
POSTER_PAGES = {1, 4, 6, 98, 99, 100}


def add_poster_slide(prs: Presentation, page: Page, total_pages: int):
    """Wide single-column layout for cover/poster/closing pages."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, C_PRIMARY)

    add_text(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.4),
        page.module or "",
        size=Pt(12),
        color=C_PRIMARY,
        bold=True,
    )
    add_text(
        slide,
        Inches(10.0),
        Inches(0.3),
        Inches(3.0),
        Inches(0.4),
        f"{page.page_type or ''}   ·   p{page.n:03d}/{total_pages}",
        size=Pt(12),
        color=C_MUTED,
        align=PP_ALIGN.RIGHT,
    )

    add_text(
        slide,
        Inches(0.5),
        Inches(0.95),
        Inches(12.5),
        Inches(1.2),
        page.title,
        size=Pt(34),
        color=C_INK,
        bold=True,
    )
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(1.5), Emu(40000), C_ACCENT)

    if page.book_pos:
        add_text(
            slide,
            Inches(0.5),
            Inches(2.15),
            Inches(12.5),
            Inches(0.35),
            f"对应书中位置：{page.book_pos}",
            size=Pt(12),
            color=C_MUTED,
        )

    # Wide bullets across full width
    if page.points:
        add_bullets(
            slide,
            Inches(0.5),
            Inches(2.8),
            Inches(12.4),
            Inches(3.8),
            page.points,
            size=Pt(20),
            color=C_INK,
        )

    _draw_takeaway_strip(slide, page)
    _set_speaker_notes(slide, page)


def _draw_takeaway_strip(slide, page: Page):
    """Slim bottom band with dynamic label per page type."""
    label = takeaway_label_for(page.page_type)
    strip_y = Inches(7.0)
    strip_h = Inches(0.5)
    add_rect(slide, 0, strip_y, SLIDE_W, strip_h, C_TAKEAWAY_BG)
    # Accent square before label
    add_rect(
        slide,
        Inches(0.5),
        strip_y + Inches(0.15),
        Inches(0.18),
        Inches(0.2),
        C_ACCENT,
    )
    add_text(
        slide,
        Inches(0.78),
        strip_y,
        Inches(1.5),
        strip_h,
        label,
        size=Pt(11),
        color=C_ACCENT,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(2.3),
        strip_y,
        Inches(10.7),
        strip_h,
        page.takeaway or "",
        size=Pt(13),
        color=C_TAKEAWAY_TEXT,
        bold=False,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_standard_slide(prs: Presentation, page: Page, total_pages: int):
    """Two-column layout: bullets left, image (or styled diagram) right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_BG)
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, C_PRIMARY)

    # Top meta row
    add_text(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(0.4),
        page.module or "",
        size=Pt(12),
        color=C_PRIMARY,
        bold=True,
    )
    add_text(
        slide,
        Inches(10.0),
        Inches(0.3),
        Inches(3.0),
        Inches(0.4),
        f"{page.page_type or ''}   ·   p{page.n:03d}/{total_pages}",
        size=Pt(12),
        color=C_MUTED,
        align=PP_ALIGN.RIGHT,
    )

    # Title + underline (1.5" accent bar)
    add_text(
        slide,
        Inches(0.5),
        Inches(0.75),
        Inches(12.5),
        Inches(1.0),
        page.title,
        size=Pt(28),
        color=C_INK,
        bold=True,
    )
    add_rect(slide, Inches(0.5), Inches(1.75), Inches(1.5), Emu(40000), C_ACCENT)

    if page.book_pos:
        add_text(
            slide,
            Inches(0.5),
            Inches(1.85),
            Inches(12.5),
            Inches(0.35),
            f"对应书中位置：{page.book_pos}",
            size=Pt(11),
            color=C_MUTED,
        )

    # Left column: bullets
    add_text(
        slide,
        Inches(0.5),
        Inches(2.35),
        Inches(2),
        Inches(0.35),
        "核心要点",
        size=Pt(13),
        color=C_PRIMARY,
        bold=True,
    )
    if page.points:
        # Smaller font for longer point lists
        bullet_size = Pt(14) if len(page.points) > 4 else Pt(16)
        add_bullets(
            slide,
            Inches(0.5),
            Inches(2.7),
            Inches(7.0),
            Inches(4.2),
            page.points,
            size=bullet_size,
            color=C_INK,
        )

    # Right column: real image OR styled text diagram
    fig_x = Inches(7.8)
    fig_y = Inches(2.3)
    fig_w = Inches(5.2)
    fig_h = Inches(4.4)

    img_rel = PAGE_IMAGES.get(page.n)
    img_path = resolve_image_path(img_rel) if img_rel else None

    if img_path:
        # White card with border
        add_rect(slide, fig_x, fig_y, fig_w, fig_h, RGBColor(0xFF, 0xFF, 0xFF), line=C_BOX_BORDER, line_w=Pt(0.75))
        # Caption strip
        add_text(
            slide,
            fig_x + Inches(0.15),
            fig_y + Inches(0.05),
            fig_w - Inches(0.3),
            Inches(0.3),
            "图源：本书配图",
            size=Pt(9),
            color=C_PRIMARY,
            bold=True,
        )
        # Image inside the card with padding
        pad = Inches(0.25)
        add_image_fitted(
            slide,
            img_path,
            int(fig_x + pad),
            int(fig_y + Inches(0.4)),
            int(fig_w - 2 * pad),
            int(fig_h - Inches(0.55)),
        )
    else:
        # Styled text diagram (configurable callout from outline.figure)
        add_rect(slide, fig_x, fig_y, fig_w, fig_h, C_BOX_BG, line=C_BOX_BORDER, line_w=Pt(1))
        add_text(
            slide,
            fig_x + Inches(0.15),
            fig_y + Inches(0.1),
            fig_w - Inches(0.3),
            Inches(0.3),
            "配图说明",
            size=Pt(11),
            color=C_PRIMARY,
            bold=True,
        )
        add_text(
            slide,
            fig_x + Inches(0.2),
            fig_y + Inches(0.5),
            fig_w - Inches(0.4),
            fig_h - Inches(0.7),
            page.figure or "（无配图建议）",
            size=Pt(13),
            color=C_INK,
            anchor=MSO_ANCHOR.TOP,
        )

    _draw_takeaway_strip(slide, page)
    _set_speaker_notes(slide, page)


def _set_speaker_notes(slide, page: Page):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = ""
    if page.speech:
        notes_tf.paragraphs[0].text = "【讲稿】" + page.speech
    if page.figure:
        para = notes_tf.add_paragraph()
        para.text = "【配图建议】" + page.figure
    if page.takeaway:
        para = notes_tf.add_paragraph()
        para.text = "【听众带走】" + page.takeaway


def build(pages: list[Page]):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = len(pages)
    add_cover_slide(prs, total)
    for page in pages:
        if page.n in POSTER_PAGES:
            add_poster_slide(prs, page, total)
        else:
            add_standard_slide(prs, page, total)

    OUTPUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PPTX))
    return prs


def main() -> int:
    md_text = INPUT_MD.read_text(encoding="utf-8")
    pages = parse_outline(md_text)
    if len(pages) != 100:
        print(f"WARN: parsed {len(pages)} pages, expected 100")

    # Image-mapping report
    have_img = 0
    missing_img = []
    for p in pages:
        rel = PAGE_IMAGES.get(p.n)
        if not rel:
            continue
        if resolve_image_path(rel):
            have_img += 1
        else:
            missing_img.append((p.n, rel))

    text_only = [p.n for p in pages if p.n not in PAGE_IMAGES]
    print(f"Parsed pages: {len(pages)}")
    print(f"  with real image : {have_img}")
    print(f"  text-only (poster/no match): {len(text_only)} -> {text_only}")
    if missing_img:
        print(f"  MISSING image files: {missing_img}")

    build(pages)
    print(f"Wrote: {OUTPUT_PPTX}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
